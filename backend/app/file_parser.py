"""
Lecture material file parser.

Extracts text content per page/slide from PDF, DOCX, and PPTX files.
"""

import io
from typing import List

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB


def parse_pdf(file_bytes: bytes) -> List[dict]:
    """Extract text per page from a PDF file using PyMuPDF (fitz).

    Returns a list of dicts with slide_number (1-indexed) and text_content.
    Returns empty list on parse failure.
    """
    if len(file_bytes) > MAX_FILE_SIZE:
        return []

    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            pages.append({
                "slide_number": i + 1,
                "text_content": text,
            })
        doc.close()
        return pages
    except Exception:
        return []


def parse_docx(file_bytes: bytes) -> List[dict]:
    """Extract text from a DOCX file, grouped into ~500-char pages.

    DOCX files don't have real page breaks, so we group paragraphs into
    chunks of approximately 500 characters each.

    Returns a list of dicts with slide_number (1-indexed) and text_content.
    Returns empty list on parse failure.
    """
    if len(file_bytes) > MAX_FILE_SIZE:
        return []

    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        if not paragraphs:
            return []

        # Group paragraphs into ~500-char pages
        pages: List[dict] = []
        current_text = ""
        page_num = 1

        for para in paragraphs:
            if current_text and len(current_text) + len(para) + 1 > 500:
                pages.append({
                    "slide_number": page_num,
                    "text_content": current_text,
                })
                page_num += 1
                current_text = para
            else:
                if current_text:
                    current_text += "\n" + para
                else:
                    current_text = para

        # Don't forget the last chunk
        if current_text:
            pages.append({
                "slide_number": page_num,
                "text_content": current_text,
            })

        return pages
    except Exception:
        return []


def parse_pptx(file_bytes: bytes) -> List[dict]:
    """Extract text per slide from a PPTX file using python-pptx.

    Extracts title and body text from all shapes on each slide.

    Returns a list of dicts with slide_number (1-indexed) and text_content.
    Returns empty list on parse failure.
    """
    if len(file_bytes) > MAX_FILE_SIZE:
        return []

    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        slides = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            slides.append({
                "slide_number": i + 1,
                "text_content": "\n".join(texts),
            })

        return slides
    except Exception:
        return []
